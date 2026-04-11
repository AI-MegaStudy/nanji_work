from pathlib import Path
import math

import pandas as pd
import matplotlib.pyplot as plt
from matplotlib import font_manager, rcParams
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
OUTPUT_DIR = BASE_DIR / "nanji_outputs"
REPORT_PATH = BASE_DIR / "nanji_weighted_ridge_modeling_report.md"
PPT_PATH = BASE_DIR / "nanji_weighted_ridge_modeling_report_presentation.pptx"
MODEL_COMPARE_PNG = OUTPUT_DIR / "nanji_operational_model_compare.png"
FEATURE_GROUP_PNG = OUTPUT_DIR / "nanji_operational_feature_groups.png"

NAVY = RGBColor(14, 34, 64)
BLUE = RGBColor(32, 91, 163)
SKY = RGBColor(83, 157, 222)
GREEN = RGBColor(20, 130, 100)
TEAL = RGBColor(45, 145, 145)
ORANGE = RGBColor(220, 126, 53)
RED = RGBColor(191, 70, 56)
TEXT = RGBColor(32, 41, 56)
MUTED = RGBColor(96, 107, 124)
LIGHT = RGBColor(240, 244, 248)
WHITE = RGBColor(255, 255, 255)


def set_korean_font():
    candidates = ["Malgun Gothic", "AppleGothic", "NanumGothic", "맑은 고딕"]
    available = {f.name for f in font_manager.fontManager.ttflist}
    for candidate in candidates:
        if candidate in available:
            rcParams["font.family"] = candidate
            return candidate
    return None


def add_textbox(slide, left, top, width, height, text="", size=20, bold=False,
                color=TEXT, font_name="Malgun Gothic", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, bullets, size=18, color=TEXT,
                level0_bold=None, font_name="Malgun Gothic", spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    level0_bold = set(level0_bold or [])
    for idx, item in enumerate(bullets):
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = text
        p.level = level
        p.space_after = Pt(6)
        p.line_spacing = spacing
        run = p.runs[0]
        run.font.name = font_name
        run.font.size = Pt(size - (1 if level else 0))
        run.font.color.rgb = color
        run.font.bold = text in level0_bold
    return box


def add_rect(slide, left, top, width, height, fill, line=None, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_title_block(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.65), Inches(0.3), Inches(12), Inches(0.5),
                title, size=28, bold=True, color=NAVY)
    if subtitle:
        add_textbox(slide, Inches(0.68), Inches(0.82), Inches(11.8), Inches(0.35),
                    subtitle, size=11, color=MUTED)
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(1.15), Inches(12.0), Inches(0.04)).fill.solid()
    bar = slide.shapes[-1]
    bar.fill.fore_color.rgb = SKY
    bar.line.color.rgb = SKY


def format_metric(v):
    return f"{v:.3f}" if isinstance(v, float) else str(v)


def create_supporting_charts():
    set_korean_font()
    weather_df = pd.read_csv(OUTPUT_DIR / "nanji_weather_only_comparison.csv")
    ordered = weather_df.sort_values("r2", ascending=True)

    fig, ax = plt.subplots(figsize=(8.4, 4.6), dpi=180)
    colors = [tuple(SKY), tuple(BLUE), tuple(TEAL), tuple(GREEN), tuple(ORANGE)]
    color_list = [
        f"#{c[0]:02x}{c[1]:02x}{c[2]:02x}"
        for c in colors[:len(ordered)]
    ]
    bars = ax.barh(ordered["model_name"], ordered["r2"], color=color_list)
    ax.set_xlim(0.70, 0.78)
    ax.set_xlabel("Test R²")
    ax.set_title("운영 후보 모델 비교")
    ax.grid(axis="x", alpha=0.2)
    for bar, val in zip(bars, ordered["r2"]):
        ax.text(val + 0.001, bar.get_y() + bar.get_height() / 2, f"{val:.3f}",
                va="center", fontsize=9)
    plt.tight_layout()
    fig.savefig(MODEL_COMPARE_PNG, bbox_inches="tight")
    plt.close(fig)

    feature_df = pd.read_csv(OUTPUT_DIR / "nanji_final_operational_feature_list.csv")
    grouped = feature_df.groupby("feature_group").size().sort_values(ascending=True)
    fig, ax = plt.subplots(figsize=(6.4, 3.8), dpi=180)
    group_colors = ["#8ab6e8", "#5aa1d6", "#3d8e96", "#2e6f95"]
    bars = ax.barh(grouped.index, grouped.values, color=group_colors[: len(grouped)])
    ax.set_title("최종 운영 feature 구성")
    ax.set_xlabel("Count")
    ax.grid(axis="x", alpha=0.2)
    for bar, val in zip(bars, grouped.values):
        ax.text(val + 0.05, bar.get_y() + bar.get_height() / 2, f"{val}",
                va="center", fontsize=9)
    plt.tight_layout()
    fig.savefig(FEATURE_GROUP_PNG, bbox_inches="tight")
    plt.close(fig)


def build_presentation():
    create_supporting_charts()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    font_name = set_korean_font() or "Malgun Gothic"

    report_title = "난지 한강공원 시간별 주차 예측 보고서"
    report_sub = "가중치 기반 Ridge 모델링 결과를 운영 관점으로 요약한 가로형 보고서"

    metrics = pd.read_csv(OUTPUT_DIR / "nanji_weather_only_comparison.csv")
    final_row = metrics.loc[metrics["model_name"] == "weather_only_extended_final"].iloc[0]
    weighted_row = metrics.loc[metrics["model_name"] == "weighted_extended"].iloc[0]
    model_df = pd.read_csv(OUTPUT_DIR / "nanji_model_metrics.csv")
    outlier_df = pd.read_csv(OUTPUT_DIR / "nanji_outlier_rule_summary.csv")
    pruned_df = pd.read_csv(OUTPUT_DIR / "nanji_feature_pruning.csv")
    removed_rows_df = pd.read_csv(OUTPUT_DIR / "nanji_feature_removed_rows.csv")
    feature_df = pd.read_csv(OUTPUT_DIR / "nanji_final_operational_feature_list.csv")

    slide = prs.slides.add_slide(blank)
    bg = add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, fill=WHITE)
    bg.line.color.rgb = WHITE
    add_rect(slide, Inches(0.65), Inches(0.8), Inches(5.9), Inches(5.6), fill=LIGHT, radius=True)
    add_rect(slide, Inches(6.9), Inches(0), Inches(6.5), prs.slide_height, fill=NAVY)
    add_textbox(slide, Inches(0.95), Inches(1.15), Inches(5.2), Inches(1.0),
                report_title, size=28, bold=True, color=NAVY, font_name=font_name)
    add_textbox(slide, Inches(0.98), Inches(2.35), Inches(5.0), Inches(1.0),
                "시간대별 점유 차량 수를 예측하고\n나중에 여유 주차공간 수로 변환할 수 있게 정리한 모델링 보고서",
                size=18, color=TEXT, font_name=font_name)
    add_bullets(
        slide, Inches(0.98), Inches(3.55), Inches(5.0), Inches(2.2),
        [
            "운영시간 기준: 06:00~23:00",
            "데이터 기간: 2022~2025, 총 34,920행",
            f"최종 운영 모델: weather_only_extended_final (R² {final_row['r2']:.3f})",
            "출처: nanji_weighted_ridge_modeling_report.md / nanji_outputs"
        ],
        size=17, color=TEXT, font_name=font_name
    )
    add_textbox(slide, Inches(7.45), Inches(1.05), Inches(4.8), Inches(0.6),
                "핵심 요약", size=24, bold=True, color=WHITE, font_name=font_name)
    summary_items = [
        "기본 패턴을 먼저 만들고, 월/시간 가중치와 날씨 변수를 더해 Ridge로 보정",
        f"단순 운영 가능한 13개 feature만 남겨도 Test RMSE {final_row['rmse']:.2f}, MAE {final_row['mae']:.2f}",
        "고혼잡 구간은 과소예측 경향이 있어 운영상 보수적 해석이 필요"
    ]
    y = 1.95
    for item in summary_items:
        add_rect(slide, Inches(7.35), Inches(y), Inches(5.2), Inches(0.95), fill=RGBColor(27, 55, 96), radius=True)
        add_textbox(slide, Inches(7.6), Inches(y + 0.12), Inches(4.8), Inches(0.72),
                    item, size=16, color=WHITE, font_name=font_name)
        y += 1.15

    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "1. 분석 목적과 데이터 범위",
                    "보고서 1~4장 내용을 발표용으로 압축")
    add_rect(slide, Inches(0.7), Inches(1.45), Inches(4.0), Inches(4.9), fill=LIGHT, radius=True)
    add_textbox(slide, Inches(0.95), Inches(1.7), Inches(3.5), Inches(0.5),
                "프로젝트 목적", size=22, bold=True, color=NAVY, font_name=font_name)
    add_bullets(slide, Inches(0.95), Inches(2.35), Inches(3.4), Inches(3.5), [
        "타깃: 특정 시각 주차장 내 추정 점유 차량 수 `estimated_active_cars`",
        "목표: 시간대별 주차 수요를 예측하고 `여유 주차공간 수`로 해석 가능하게 구성",
        "현장 실측 시간 로그가 아니라 일별 원본을 시간대별로 확장한 통합 데이터셋 활용"
    ], size=17, font_name=font_name)
    add_rect(slide, Inches(4.95), Inches(1.45), Inches(3.3), Inches(4.9), fill=WHITE, line=RGBColor(214, 223, 232), radius=True)
    add_textbox(slide, Inches(5.2), Inches(1.7), Inches(2.7), Inches(0.5),
                "운영 기준", size=22, bold=True, color=NAVY, font_name=font_name)
    add_bullets(slide, Inches(5.2), Inches(2.35), Inches(2.7), Inches(3.5), [
        "운영시간은 `06~23시`",
        "`00~05시`는 데이터가 있어도 예측값을 0으로 유지",
        "학습 구간 분리: train 2022~2023 / valid 2024 / test 2025"
    ], size=16, font_name=font_name)
    add_rect(slide, Inches(8.5), Inches(1.45), Inches(4.15), Inches(4.9), fill=WHITE, line=RGBColor(214, 223, 232), radius=True)
    add_textbox(slide, Inches(8.75), Inches(1.7), Inches(3.7), Inches(0.5),
                "데이터 현황", size=22, bold=True, color=NAVY, font_name=font_name)
    add_bullets(slide, Inches(8.75), Inches(2.35), Inches(3.5), Inches(3.6), [
        "총 34,920행, datetime 중복 0, 타깃 결측 0",
        "기간: 2022-01-01 00:00 ~ 2025-12-31 23:00",
        "날씨 커버리지 100%, 지하철 92.4%, 버스 72.0%",
        "자전거 33.2%, 문화행사 18.8%로 일부 변수는 보조 설명용"
    ], size=16, font_name=font_name)

    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "2. 모델링 구조",
                    "기본 패턴을 만든 뒤 가중치와 Ridge 보정을 얹는 해석형 구조")
    steps = [
        ("1", "기본 패턴", "day_type x hour 평균으로\nbase_value 생성"),
        ("2", "월 보정", "month_weight로\n계절 수준 차이 반영"),
        ("3", "시간 보정", "hour_weight로\n운영시간 내 혼잡 편차 반영"),
        ("4", "Ridge 보정", "주기 변수 + 휴일 + 날씨로\n잔차를 줄이는 최종 단계"),
    ]
    x_positions = [0.9, 3.25, 5.6, 7.95]
    for (num, head, body), x in zip(steps, x_positions):
        add_rect(slide, Inches(x), Inches(2.1), Inches(2.0), Inches(2.3), fill=LIGHT, radius=True)
        add_textbox(slide, Inches(x + 0.12), Inches(2.24), Inches(0.45), Inches(0.4),
                    num, size=18, bold=True, color=WHITE, font_name=font_name, align=PP_ALIGN.CENTER)
        pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.12), Inches(2.22), Inches(0.42), Inches(0.42))
        pill.fill.solid()
        pill.fill.fore_color.rgb = BLUE
        pill.line.color.rgb = BLUE
        add_textbox(slide, Inches(x + 0.65), Inches(2.22), Inches(1.1), Inches(0.35),
                    head, size=18, bold=True, color=NAVY, font_name=font_name)
        add_textbox(slide, Inches(x + 0.15), Inches(2.95), Inches(1.7), Inches(0.95),
                    body, size=15, color=TEXT, font_name=font_name, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.95), Inches(5.0), Inches(11.6), Inches(1.0),
                "핵심은 블랙박스 예측이 아니라, 사람이 이해 가능한 패턴을 먼저 세우고\n머신러닝은 설명이 안 되는 차이만 보정하도록 제한한 점입니다.",
                size=18, color=TEXT, font_name=font_name, align=PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "3. 전처리와 데이터 품질 점검",
                    "분포, 이상값, 다중공선성을 먼저 확인한 뒤 모델링")
    add_rect(slide, Inches(0.7), Inches(1.45), Inches(4.2), Inches(4.8), fill=LIGHT, radius=True)
    add_textbox(slide, Inches(0.95), Inches(1.7), Inches(3.8), Inches(0.45),
                "타깃 이상값 점검", size=22, bold=True, color=NAVY, font_name=font_name)
    outlier = outlier_df.iloc[0]
    add_bullets(slide, Inches(0.95), Inches(2.3), Inches(3.7), Inches(3.2), [
        f"IQR 상한 후보 {int(outlier['candidate_count']):,}건 모두 특이치로 유지",
        f"실제 제거된 target outlier는 {int(outlier['outlier_count'])}건",
        "즉 큰 값은 드물지만 끊어진 오류라기보다 실제 혼잡 이벤트로 해석"
    ], size=17, font_name=font_name)
    add_rect(slide, Inches(5.05), Inches(1.45), Inches(3.35), Inches(4.8), fill=WHITE, line=RGBColor(214, 223, 232), radius=True)
    add_textbox(slide, Inches(5.3), Inches(1.7), Inches(2.8), Inches(0.45),
                "입력 변수 정리", size=22, bold=True, color=NAVY, font_name=font_name)
    add_bullets(slide, Inches(5.3), Inches(2.3), Inches(2.8), Inches(3.5), [
        f"상관계수 0.9 이상 기준으로 {len(pruned_df):,}개 중복 변수 정리",
        f"실제 행 제거는 자전거 관련 극단치 {len(removed_rows_df):,}건만 반영",
        "로그 변환은 검토만 하고 최종 구조는 원해석 가능성을 유지"
    ], size=16, font_name=font_name)
    add_rect(slide, Inches(8.55), Inches(1.45), Inches(4.1), Inches(4.8), fill=WHITE, line=RGBColor(214, 223, 232), radius=True)
    add_textbox(slide, Inches(8.8), Inches(1.7), Inches(3.5), Inches(0.45),
                "해석 포인트", size=22, bold=True, color=NAVY, font_name=font_name)
    add_bullets(slide, Inches(8.8), Inches(2.3), Inches(3.45), Inches(3.6), [
        "분포는 작은 값이 많고 큰 값은 드문 long-tail 구조",
        "선형 모델은 일반 구간에 안정적이지만, 피크 구간은 놓칠 수 있음",
        "따라서 전체 정확도와 운영 단순성 사이의 균형이 중요"
    ], size=16, font_name=font_name)

    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "4. 패턴 가중치 시각화",
                    "월별 수준 차이와 운영시간 내 혼잡 편차를 분리해 해석")
    slide.shapes.add_picture(str(OUTPUT_DIR / "nanji_month_weights.png"), Inches(0.8), Inches(1.6), width=Inches(5.7))
    slide.shapes.add_picture(str(OUTPUT_DIR / "nanji_hour_weights.png"), Inches(6.8), Inches(1.6), width=Inches(5.7))
    add_textbox(slide, Inches(0.95), Inches(5.95), Inches(5.4), Inches(0.65),
                "월 가중치: 4~6월, 9~10월이 높고 1~2월·12월이 낮음", size=16, color=TEXT, font_name=font_name, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.95), Inches(5.95), Inches(5.4), Inches(0.65),
                "시간 가중치: 13~16시에 혼잡이 집중되고 06~10시는 상대적으로 낮음", size=16, color=TEXT, font_name=font_name, align=PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "5. 모델 비교와 최종 운영안",
                    "정확도뿐 아니라 미래 시점에서 확보 가능한 입력인지 함께 고려")
    slide.shapes.add_picture(str(MODEL_COMPARE_PNG), Inches(0.8), Inches(1.55), width=Inches(6.0))
    table = slide.shapes.add_table(6, 5, Inches(7.0), Inches(1.7), Inches(5.55), Inches(3.0)).table
    headers = ["모델", "alpha", "RMSE", "MAE", "R²"]
    rows = metrics.sort_values("r2", ascending=False).reset_index(drop=True)
    for col, head in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = head
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.color.rgb = WHITE
                r.font.name = font_name
                r.font.size = Pt(11)
    for ridx in range(len(rows)):
        row = rows.iloc[ridx]
        values = [row["model_name"], f"{int(row['alpha'])}", f"{row['rmse']:.2f}", f"{row['mae']:.2f}", f"{row['r2']:.3f}"]
        for cidx, value in enumerate(values):
            cell = table.cell(ridx + 1, cidx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if row["model_name"] != "weather_only_extended_final" else RGBColor(220, 240, 234)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = font_name
                    r.font.size = Pt(10.5)
                    r.font.bold = row["model_name"] == "weather_only_extended_final"
                    r.font.color.rgb = TEXT
    add_bullets(slide, Inches(7.05), Inches(5.0), Inches(5.3), Inches(1.2), [
        f"최종 운영안: `weather_only_extended_final` (Test R² {final_row['r2']:.3f})",
        f"비교 기준 `weighted_extended` 대비 RMSE를 {weighted_row['rmse'] - final_row['rmse']:.2f} 개선",
        "교통·행사 변수 없이도 날씨 예보와 패턴 정보만으로 운영 가능한 구조"
    ], size=15, font_name=font_name)

    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "6. 최종 운영 feature",
                    "총 13개 변수로 축소해 성능과 운영 가능성을 함께 확보")
    slide.shapes.add_picture(str(FEATURE_GROUP_PNG), Inches(0.8), Inches(1.7), width=Inches(4.8))
    add_rect(slide, Inches(5.95), Inches(1.55), Inches(6.2), Inches(4.8), fill=WHITE, line=RGBColor(214, 223, 232), radius=True)
    add_textbox(slide, Inches(6.2), Inches(1.8), Inches(5.5), Inches(0.45),
                "feature 구성", size=22, bold=True, color=NAVY, font_name=font_name)
    features = [
        "패턴: pattern_prior, month_weight, hour_weight, day_type_offday",
        "주기: hour_sin, hour_cos, month_sin, month_cos",
        "달력: is_holiday",
        "날씨: temperature_2m, relative_humidity_2m, weather_code, wind_gusts_10m",
    ]
    add_bullets(slide, Inches(6.2), Inches(2.35), Inches(5.4), Inches(2.7), features, size=16, font_name=font_name)
    add_textbox(slide, Inches(6.2), Inches(5.3), Inches(5.3), Inches(0.7),
                "미래 예측 시점에 확보 가능한 입력만 남겨 실제 운영 흐름과 연결하기 쉽게 만든 구성이 핵심입니다.",
                size=15, color=TEXT, font_name=font_name)

    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "7. 2025년 테스트셋 결과 해석",
                    "월별 흐름은 따라가지만, 매우 높은 피크는 보수적으로 보는 것이 안전")
    slide.shapes.add_picture(str(OUTPUT_DIR / "nanji_2025_monthly_compare.png"), Inches(0.8), Inches(1.55), width=Inches(7.0))
    add_rect(slide, Inches(8.0), Inches(1.55), Inches(4.6), Inches(4.85), fill=LIGHT, radius=True)
    add_textbox(slide, Inches(8.25), Inches(1.8), Inches(4.0), Inches(0.45),
                "운영 해석", size=22, bold=True, color=NAVY, font_name=font_name)
    add_bullets(slide, Inches(8.25), Inches(2.35), Inches(3.95), Inches(3.5), [
        "100 미만과 100~300 구간은 비교적 안정적으로 설명",
        "300~500 구간부터 오차가 커지고, 500 이상은 과소예측 경향이 뚜렷",
        "즉 일반 운영 수요 추세 예측에는 유용하지만, 극단적 피크는 별도 관리가 필요",
        "보고서 기준 최종 해석: 간단한 Ridge 운영안 + 고혼잡 구간에 대한 보수적 경보 체계"
    ], size=16, font_name=font_name)

    slide = prs.slides.add_slide(blank)
    add_title_block(slide, "8. 활용 방식과 제안",
                    "예측 점유 차량 수를 여유 주차공간으로 변환해 바로 운영지표로 사용")
    add_rect(slide, Inches(0.8), Inches(1.75), Inches(12.0), Inches(1.2), fill=LIGHT, radius=True)
    add_textbox(slide, Inches(1.1), Inches(2.05), Inches(11.3), Inches(0.55),
                "available_spaces = max(total_capacity - predicted_estimated_active_cars, 0)",
                size=24, bold=True, color=NAVY, font_name=font_name, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(1.0), Inches(3.45), Inches(3.45), Inches(2.2), fill=WHITE, line=RGBColor(214, 223, 232), radius=True)
    add_rect(slide, Inches(4.95), Inches(3.45), Inches(3.45), Inches(2.2), fill=WHITE, line=RGBColor(214, 223, 232), radius=True)
    add_rect(slide, Inches(8.9), Inches(3.45), Inches(3.45), Inches(2.2), fill=WHITE, line=RGBColor(214, 223, 232), radius=True)
    add_textbox(slide, Inches(1.2), Inches(3.75), Inches(3.0), Inches(1.2),
                "입력\n시간대 패턴 + 휴일 + 날씨 예보", size=19, bold=True, color=NAVY, font_name=font_name, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(5.15), Inches(3.75), Inches(3.0), Inches(1.2),
                "모델 출력\n예상 점유 차량 수", size=19, bold=True, color=NAVY, font_name=font_name, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(9.1), Inches(3.72), Inches(3.0), Inches(1.2),
                "운영 지표\n여유 주차공간 수 / 여유율", size=19, bold=True, color=NAVY, font_name=font_name, align=PP_ALIGN.CENTER)
    add_bullets(slide, Inches(0.95), Inches(6.0), Inches(11.8), Inches(0.7), [
        "권장 운영 방식: 1시간·2시간·하루 단위 수요 추세 안내에는 적극 활용, 고혼잡 피크는 보수적 해석과 현장 모니터링 병행"
    ], size=15, font_name=font_name)

    prs.save(PPT_PATH)
    return PPT_PATH


if __name__ == "__main__":
    path = build_presentation()
    print(path)
