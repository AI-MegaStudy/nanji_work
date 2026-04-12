from pathlib import Path

import pandas as pd
import matplotlib.pyplot as plt
from matplotlib import font_manager, rcParams
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
SOURCE_DIR = ROOT / "hmw" / "Note"
OUTPUT_DIR = Path(__file__).resolve().parent
SOURCE_OUT = SOURCE_DIR / "nanji_outputs"
PPT_PATH = OUTPUT_DIR / "난지_주차예측_보고서_업데이트_가로형.pptx"
MODEL_COMPARE_IMG = OUTPUT_DIR / "nanji_operational_model_compare_update.png"
FEATURE_GROUP_IMG = OUTPUT_DIR / "nanji_operational_feature_groups_update.png"

NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(29, 78, 216)
LINK = RGBColor(30, 58, 138)
TEXT = RGBColor(17, 24, 39)
MUTED = RGBColor(71, 85, 105)
BORDER = RGBColor(148, 163, 184)
PANEL = RGBColor(255, 255, 255)
SOFT = RGBColor(239, 246, 255)
TH_BG = RGBColor(224, 236, 255)
WHITE = RGBColor(255, 255, 255)


def set_font():
    candidates = ["Malgun Gothic", "Apple SD Gothic Neo", "AppleGothic", "NanumGothic", "맑은 고딕"]
    available = {f.name for f in font_manager.fontManager.ttflist}
    for name in candidates:
        if name in available:
            rcParams["font.family"] = name
            return name
    return "Malgun Gothic"


FONT_NAME = set_font()


def rgb_tuple(color):
    return tuple(color)


def add_box(slide, left, top, width, height, fill, line=None, rounded=True):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_text(slide, left, top, width, height, text, size=20, bold=False,
             color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.alignment = align
    p.text = text
    run = p.runs[0]
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(5)
        p.line_spacing = 1.15
        run = p.runs[0]
        run.font.name = FONT_NAME
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_header(slide, title, subtitle=None):
    add_text(slide, Inches(0.5), Inches(0.28), Inches(10.2), Inches(0.45), title, size=23, bold=True, color=NAVY)
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.5), Inches(0.76), Inches(10.2), Inches(0.035))
    line = slide.shapes[-1]
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.color.rgb = BLUE
    if subtitle:
        add_text(slide, Inches(0.52), Inches(0.82), Inches(10.1), Inches(0.24), subtitle, size=9, color=MUTED)


def make_supporting_images():
    weather = pd.read_csv(SOURCE_OUT / "nanji_weather_only_comparison.csv").sort_values("r2", ascending=True)
    colors = ["#94a3b8", "#60a5fa", "#2563eb", "#1e3a8a", "#0f172a"]
    fig, ax = plt.subplots(figsize=(7.2, 3.7), dpi=180)
    bars = ax.barh(weather["model_name"], weather["r2"], color=colors[: len(weather)])
    ax.set_xlim(0.72, 0.77)
    ax.set_xlabel("Test R²")
    ax.set_title("운영 후보 모델 비교")
    ax.grid(axis="x", alpha=0.2)
    for bar, val in zip(bars, weather["r2"]):
        ax.text(val + 0.0006, bar.get_y() + bar.get_height() / 2, f"{val:.3f}", va="center", fontsize=8)
    plt.tight_layout()
    fig.savefig(MODEL_COMPARE_IMG, bbox_inches="tight")
    plt.close(fig)

    features = pd.read_csv(SOURCE_OUT / "nanji_final_operational_feature_list.csv")
    group_counts = features.groupby("feature_group").size().sort_values(ascending=True)
    fig, ax = plt.subplots(figsize=(5.2, 3.4), dpi=180)
    bars = ax.barh(group_counts.index, group_counts.values, color=["#cbd5e1", "#93c5fd", "#60a5fa", "#1d4ed8"])
    ax.set_xlabel("Count")
    ax.set_title("최종 운영 feature 구성")
    ax.grid(axis="x", alpha=0.2)
    for bar, val in zip(bars, group_counts.values):
        ax.text(val + 0.03, bar.get_y() + bar.get_height() / 2, str(val), va="center", fontsize=8)
    plt.tight_layout()
    fig.savefig(FEATURE_GROUP_IMG, bbox_inches="tight")
    plt.close(fig)


def style_table_cell(cell, fill, color, size=10, bold=False):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            r.font.name = FONT_NAME
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color


def build_ppt():
    make_supporting_images()

    metrics = pd.read_csv(SOURCE_OUT / "nanji_weather_only_comparison.csv")
    final_row = metrics.loc[metrics["model_name"] == "weather_only_extended_final"].iloc[0]
    weighted_row = metrics.loc[metrics["model_name"] == "weighted_extended"].iloc[0]
    outlier = pd.read_csv(SOURCE_OUT / "nanji_outlier_rule_summary.csv").iloc[0]
    removed_rows = pd.read_csv(SOURCE_OUT / "nanji_feature_removed_rows.csv")
    feature_pruning = pd.read_csv(SOURCE_OUT / "nanji_feature_pruning.csv")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_box(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE, WHITE, rounded=False)
    add_text(slide, Inches(0.55), Inches(0.62), Inches(8.7), Inches(0.8),
             "난지 한강공원 시간별 주차 예측 보고서", size=28, bold=True, color=NAVY)
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.38), Inches(8.9), Inches(0.045))
    line = slide.shapes[-1]
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.color.rgb = BLUE
    add_text(slide, Inches(0.58), Inches(1.62), Inches(8.4), Inches(0.95),
             "A4 가로형 보고서 스타일 업데이트 버전\n가중치 기반 Ridge 모델링 결과를 운영 관점으로 요약",
             size=18, color=TEXT)
    add_box(slide, Inches(0.58), Inches(2.7), Inches(5.2), Inches(2.0), SOFT, BLUE)
    add_text(slide, Inches(0.82), Inches(2.92), Inches(4.7), Inches(0.35), "핵심 결론", size=18, bold=True, color=NAVY)
    add_bullets(slide, Inches(0.82), Inches(3.35), Inches(4.65), Inches(1.15), [
        "최종 운영안은 `weather_only_extended_final`",
        f"Test RMSE {final_row['rmse']:.2f}, MAE {final_row['mae']:.2f}, R² {final_row['r2']:.3f}",
        "날씨 예보와 패턴 정보만으로 운영 가능"
    ], size=15)
    add_box(slide, Inches(6.2), Inches(2.7), Inches(4.85), Inches(2.0), PANEL, BORDER)
    add_text(slide, Inches(6.45), Inches(2.92), Inches(4.2), Inches(0.35), "적용한 스타일 방향", size=18, bold=True, color=NAVY)
    add_bullets(slide, Inches(6.45), Inches(3.35), Inches(4.0), Inches(1.05), [
        "짙은 남색 제목 + 파란 구분선",
        "밝은 카드형 레이아웃과 테이블 강조",
        "A4 landscape 비율에 맞춘 보고서형 슬라이드"
    ], size=15)
    add_text(slide, Inches(0.58), Inches(7.45), Inches(5.5), Inches(0.25),
             "출력 위치: yeeun/0411", size=10, color=MUTED)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "1. 분석 목적과 데이터 범위", "보고서 1~4장 핵심 요약")
    add_box(slide, Inches(0.55), Inches(1.25), Inches(3.35), Inches(5.8), PANEL, BORDER)
    add_text(slide, Inches(0.8), Inches(1.5), Inches(2.8), Inches(0.3), "프로젝트 목적", size=18, bold=True, color=NAVY)
    add_bullets(slide, Inches(0.8), Inches(1.9), Inches(2.8), Inches(4.6), [
        "타깃은 특정 시각 주차장 내 추정 점유 차량 수 `estimated_active_cars`",
        "시간대별 수요를 예측하고 나중에 여유 주차공간 수로 변환 가능하게 구성",
        "일별 원본을 시간대별로 확장한 통합 데이터셋을 사용"
    ], size=15)
    add_box(slide, Inches(4.05), Inches(1.25), Inches(3.0), Inches(5.8), PANEL, BORDER)
    add_text(slide, Inches(4.3), Inches(1.5), Inches(2.4), Inches(0.3), "운영 기준", size=18, bold=True, color=NAVY)
    add_bullets(slide, Inches(4.3), Inches(1.9), Inches(2.35), Inches(4.5), [
        "운영시간은 `06~23시`",
        "`00~05시`는 예측값을 0으로 유지",
        "train: 2022~2023",
        "valid: 2024",
        "test: 2025"
    ], size=15)
    add_box(slide, Inches(7.2), Inches(1.25), Inches(3.95), Inches(5.8), PANEL, BORDER)
    add_text(slide, Inches(7.45), Inches(1.5), Inches(3.3), Inches(0.3), "데이터 현황", size=18, bold=True, color=NAVY)
    add_bullets(slide, Inches(7.45), Inches(1.9), Inches(3.2), Inches(4.8), [
        "총 34,920행, datetime 중복 0, 타깃 결측 0",
        "기간: 2022-01-01 00:00 ~ 2025-12-31 23:00",
        "날씨 feature coverage 100%",
        "지하철 92.4%, 버스 72.0%, 자전거 33.2%, 문화행사 18.8%"
    ], size=15)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "2. 모델링 구조", "기본 패턴을 먼저 세우고 Ridge는 남는 오차만 보정")
    steps = [
        ("기본 패턴", "day_type x hour 평균으로\nbase_value 생성"),
        ("월 보정", "month_weight로\n계절 수준 차이 반영"),
        ("시간 보정", "hour_weight로\n운영시간 내 편차 반영"),
        ("Ridge", "주기 변수, 휴일,\n날씨로 잔차 보정"),
    ]
    for idx, (title, body) in enumerate(steps):
        left = 0.6 + idx * 2.75
        add_box(slide, Inches(left), Inches(2.0), Inches(2.35), Inches(2.4), SOFT if idx % 2 == 0 else PANEL, BORDER)
        add_text(slide, Inches(left + 0.18), Inches(2.22), Inches(1.95), Inches(0.35), title, size=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(left + 0.18), Inches(2.95), Inches(1.95), Inches(0.85), body, size=14, color=TEXT, align=PP_ALIGN.CENTER)
    add_box(slide, Inches(0.8), Inches(5.05), Inches(9.9), Inches(1.35), SOFT, BLUE)
    add_text(slide, Inches(1.0), Inches(5.35), Inches(9.5), Inches(0.7),
             "핵심은 블랙박스 예측보다 해석 가능한 구조를 유지하는 것입니다. 사람도 이해할 수 있는 패턴을 먼저 만들고, 머신러닝은 설명이 안 되는 차이만 줄이도록 설계했습니다.",
             size=16, color=TEXT, align=PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "3. 전처리와 품질 점검", "분포, 이상값, 다중공선성을 먼저 정리")
    add_box(slide, Inches(0.55), Inches(1.25), Inches(3.35), Inches(5.9), PANEL, BORDER)
    add_text(slide, Inches(0.8), Inches(1.5), Inches(2.8), Inches(0.3), "타깃 이상값 점검", size=18, bold=True, color=NAVY)
    add_bullets(slide, Inches(0.8), Inches(1.9), Inches(2.7), Inches(4.7), [
        f"IQR 상한 후보 {int(outlier['candidate_count']):,}건 모두 특이치로 유지",
        f"실제 제거된 target outlier는 {int(outlier['outlier_count'])}건",
        "즉 높은 값은 드물지만 오류보다 실제 혼잡 이벤트로 해석"
    ], size=15)
    add_box(slide, Inches(4.05), Inches(1.25), Inches(3.0), Inches(5.9), PANEL, BORDER)
    add_text(slide, Inches(4.3), Inches(1.5), Inches(2.4), Inches(0.3), "입력 변수 정리", size=18, bold=True, color=NAVY)
    add_bullets(slide, Inches(4.3), Inches(1.9), Inches(2.35), Inches(4.8), [
        f"상관계수 0.9 이상 기준으로 {len(feature_pruning):,}개 중복 관계 정리",
        f"실제 행 제거는 자전거 관련 극단치 {len(removed_rows):,}건만 반영",
        "로그 변환은 검토만 하고 최종 해석 가능성 유지"
    ], size=15)
    add_box(slide, Inches(7.2), Inches(1.25), Inches(3.95), Inches(5.9), PANEL, BORDER)
    add_text(slide, Inches(7.45), Inches(1.5), Inches(3.3), Inches(0.3), "해석 포인트", size=18, bold=True, color=NAVY)
    add_bullets(slide, Inches(7.45), Inches(1.9), Inches(3.15), Inches(4.9), [
        "분포는 작은 값이 많고 큰 값은 드문 long-tail 구조",
        "선형 모델은 일반 구간에 안정적이지만 피크 구간에는 약함",
        "그래서 전체 정확도와 운영 단순성의 균형이 중요"
    ], size=15)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "4. 패턴 가중치 시각화", "계절 수준 차이와 시간대 편차를 분리해 해석")
    slide.shapes.add_picture(str(SOURCE_OUT / "nanji_month_weights.png"), Inches(0.6), Inches(1.35), width=Inches(4.95))
    slide.shapes.add_picture(str(SOURCE_OUT / "nanji_hour_weights.png"), Inches(5.95), Inches(1.35), width=Inches(4.95))
    add_text(slide, Inches(0.75), Inches(6.25), Inches(4.7), Inches(0.45),
             "월 가중치: 4~6월, 9~10월이 높고 1~2월·12월이 낮음", size=14, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(6.1), Inches(6.25), Inches(4.6), Inches(0.45),
             "시간 가중치: 13~16시 중심의 혼잡 집중 패턴 확인", size=14, color=TEXT, align=PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "5. 모델 비교와 최종 운영안", "정확도뿐 아니라 미래 시점에서 확보 가능한 입력인지 함께 고려")
    slide.shapes.add_picture(str(MODEL_COMPARE_IMG), Inches(0.65), Inches(1.45), width=Inches(5.25))
    table = slide.shapes.add_table(6, 5, Inches(6.15), Inches(1.52), Inches(4.8), Inches(2.85)).table
    headers = ["모델", "alpha", "RMSE", "MAE", "R²"]
    ranked = metrics.sort_values("r2", ascending=False).reset_index(drop=True)
    for c, header in enumerate(headers):
        table.cell(0, c).text = header
        style_table_cell(table.cell(0, c), TH_BG, NAVY, size=10, bold=True)
    for r in range(len(ranked)):
        row = ranked.iloc[r]
        values = [row["model_name"], f"{int(row['alpha'])}", f"{row['rmse']:.2f}", f"{row['mae']:.2f}", f"{row['r2']:.3f}"]
        for c, value in enumerate(values):
            table.cell(r + 1, c).text = value
            fill = SOFT if row["model_name"] == "weather_only_extended_final" else WHITE
            style_table_cell(table.cell(r + 1, c), fill, TEXT, size=9.5, bold=row["model_name"] == "weather_only_extended_final")
    add_box(slide, Inches(6.15), Inches(4.65), Inches(4.8), Inches(1.8), SOFT, BLUE)
    add_bullets(slide, Inches(6.35), Inches(4.92), Inches(4.35), Inches(1.2), [
        f"최종 운영안: weather_only_extended_final",
        f"weighted_extended 대비 RMSE {weighted_row['rmse'] - final_row['rmse']:.2f} 개선",
        f"최종 성능: R² {final_row['r2']:.3f}, MAE {final_row['mae']:.2f}"
    ], size=14)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "6. 최종 운영 feature", "총 13개 변수로 축소해 운영 가능성과 설명력을 함께 확보")
    slide.shapes.add_picture(str(FEATURE_GROUP_IMG), Inches(0.7), Inches(1.5), width=Inches(4.7))
    add_box(slide, Inches(5.75), Inches(1.5), Inches(5.0), Inches(4.95), PANEL, BORDER)
    add_text(slide, Inches(5.98), Inches(1.78), Inches(4.4), Inches(0.32), "feature 구성", size=18, bold=True, color=NAVY)
    add_bullets(slide, Inches(5.98), Inches(2.2), Inches(4.35), Inches(3.8), [
        "패턴: pattern_prior, month_weight, hour_weight, day_type_offday",
        "주기: hour_sin, hour_cos, month_sin, month_cos",
        "달력: is_holiday",
        "날씨: temperature_2m, relative_humidity_2m, weather_code, wind_gusts_10m",
        "미래 예측 시점에 확보 가능한 입력만 남겨 실제 운영 흐름과 연결"
    ], size=15)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "7. 테스트 결과 해석과 활용 방식", "월별 추세 예측에는 유용하지만 고혼잡 피크는 보수적으로 해석")
    slide.shapes.add_picture(str(SOURCE_OUT / "nanji_2025_monthly_compare.png"), Inches(0.65), Inches(1.38), width=Inches(5.35))
    add_box(slide, Inches(6.2), Inches(1.38), Inches(4.75), Inches(2.45), PANEL, BORDER)
    add_text(slide, Inches(6.45), Inches(1.65), Inches(4.2), Inches(0.3), "테스트셋 해석", size=18, bold=True, color=NAVY)
    add_bullets(slide, Inches(6.45), Inches(2.05), Inches(4.1), Inches(1.45), [
        "100 미만과 100~300 구간은 비교적 안정적으로 설명",
        "300~500부터 오차가 커지고, 500 이상은 과소예측 경향이 뚜렷"
    ], size=14)
    add_box(slide, Inches(6.2), Inches(4.0), Inches(4.75), Inches(2.45), SOFT, BLUE)
    add_text(slide, Inches(6.45), Inches(4.27), Inches(4.1), Inches(0.3), "운영 적용 식", size=18, bold=True, color=NAVY)
    add_text(slide, Inches(6.45), Inches(4.7), Inches(4.05), Inches(0.45),
             "available_spaces = max(total_capacity - predicted_estimated_active_cars, 0)",
             size=14, bold=True, color=LINK, align=PP_ALIGN.CENTER)
    add_bullets(slide, Inches(6.45), Inches(5.28), Inches(4.05), Inches(0.9), [
        "1시간·2시간·하루 단위 수요 추세 안내에 활용",
        "극단적 피크는 현장 모니터링과 병행 권장"
    ], size=14)

    prs.save(PPT_PATH)
    return PPT_PATH


if __name__ == "__main__":
    output = build_ppt()
    print(output)
