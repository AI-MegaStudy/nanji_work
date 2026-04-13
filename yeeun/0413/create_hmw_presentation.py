from pathlib import Path

import matplotlib.pyplot as plt
import pandas as pd
from matplotlib import font_manager, rcParams
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
ROOT_DIR = BASE_DIR.parents[1]
SOURCE_DIR = ROOT_DIR / "hmw" / "Note"
OUT_DIR = SOURCE_DIR / "nanji_outputs_change"

PPT_PATH = BASE_DIR / "난지_ML_change_보고서_가로형.pptx"
MODEL_COMPARE_IMG = BASE_DIR / "nanji_ml_change_model_compare.png"
FEATURE_GROUP_IMG = BASE_DIR / "nanji_ml_change_feature_groups.png"
WEIGHT_IMG = BASE_DIR / "nanji_ml_change_weights.png"
MONTHLY_TREND_IMG = BASE_DIR / "nanji_ml_change_monthly_trend.png"
IMPORTANCE_IMG = BASE_DIR / "nanji_ml_change_importance.png"

NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
SKY = RGBColor(219, 234, 254)
TEAL = RGBColor(14, 116, 144)
GREEN = RGBColor(22, 163, 74)
ORANGE = RGBColor(234, 88, 12)
RED = RGBColor(220, 38, 38)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(100, 116, 139)
BORDER = RGBColor(203, 213, 225)
LIGHT = RGBColor(248, 250, 252)
SOFT = RGBColor(239, 246, 255)


def set_font():
    candidates = ["Malgun Gothic", "Apple SD Gothic Neo", "AppleGothic", "NanumGothic"]
    available = {f.name for f in font_manager.fontManager.ttflist}
    for candidate in candidates:
        if candidate in available:
            rcParams["font.family"] = candidate
            rcParams["axes.unicode_minus"] = False
            return candidate
    rcParams["axes.unicode_minus"] = False
    return "Malgun Gothic"


FONT_NAME = set_font()


def add_box(slide, left, top, width, height, fill, line=None, rounded=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_text(slide, left, top, width, height, text, size=20, bold=False,
             color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(5)
        p.line_spacing = 1.15
        run = p.runs[0]
        run.font.name = FONT_NAME
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_header(slide, title, subtitle=""):
    add_text(slide, Inches(0.55), Inches(0.3), Inches(12.0), Inches(0.45), title, size=24, bold=True, color=NAVY)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.82), Inches(12.0), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.color.rgb = BLUE
    if subtitle:
        add_text(slide, Inches(0.57), Inches(0.9), Inches(11.8), Inches(0.25), subtitle, size=10, color=MUTED)


def style_cell(cell, fill, font_color, size=10, bold=False):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = FONT_NAME
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = font_color


def make_charts():
    compare = pd.read_csv(OUT_DIR / "nanji_weather_only_comparison.csv").sort_values("r2", ascending=True)
    fig, ax = plt.subplots(figsize=(7.0, 3.8), dpi=180)
    colors = ["#cbd5e1", "#93c5fd", "#60a5fa", "#2563eb", "#1d4ed8"]
    bars = ax.barh(compare["model_name"], compare["r2"], color=colors[: len(compare)])
    ax.set_xlim(0.58, 0.70)
    ax.set_xlabel("Test R²")
    ax.set_title("ML Change 모델 성능 비교")
    ax.grid(axis="x", alpha=0.2)
    for bar, val in zip(bars, compare["r2"]):
        ax.text(val + 0.001, bar.get_y() + bar.get_height() / 2, f"{val:.3f}", va="center", fontsize=8)
    plt.tight_layout()
    fig.savefig(MODEL_COMPARE_IMG, bbox_inches="tight")
    plt.close(fig)

    features = pd.read_csv(OUT_DIR / "nanji_final_operational_feature_list.csv")
    group_counts = features.groupby("feature_group").size().sort_values(ascending=True)
    fig, ax = plt.subplots(figsize=(5.6, 3.6), dpi=180)
    bars = ax.barh(group_counts.index, group_counts.values, color=["#dbeafe", "#93c5fd", "#38bdf8", "#0ea5e9"])
    ax.set_xlabel("Count")
    ax.set_title("운영형 최종 feature 구성")
    ax.grid(axis="x", alpha=0.2)
    for bar, val in zip(bars, group_counts.values):
        ax.text(val + 0.03, bar.get_y() + bar.get_height() / 2, str(val), va="center", fontsize=8)
    plt.tight_layout()
    fig.savefig(FEATURE_GROUP_IMG, bbox_inches="tight")
    plt.close(fig)

    weights = pd.read_csv(OUT_DIR / "nanji_weight_table.csv")
    month_df = weights.loc[weights["weight_type"] == "month_weight"].copy()
    month_df["key"] = month_df["key"].astype(int)
    hour_df = weights.loc[weights["weight_type"] == "hour_weight_raw"].copy()
    hour_df["key"] = hour_df["key"].astype(int)

    fig, axes = plt.subplots(1, 2, figsize=(9.5, 3.6), dpi=180)
    axes[0].bar(month_df["key"], month_df["value"], color="#2563eb")
    axes[0].set_title("월 가중치")
    axes[0].set_xlabel("Month")
    axes[0].grid(axis="y", alpha=0.2)
    axes[0].set_xticks(range(1, 13))

    axes[1].plot(hour_df["key"], hour_df["value"], color="#ea580c", linewidth=2.5, marker="o", markersize=4)
    axes[1].set_title("시간 가중치")
    axes[1].set_xlabel("Hour")
    axes[1].grid(axis="y", alpha=0.2)
    axes[1].set_xticks(sorted(hour_df["key"].tolist()))

    plt.tight_layout()
    fig.savefig(WEIGHT_IMG, bbox_inches="tight")
    plt.close(fig)

    test_df = pd.read_csv(OUT_DIR / "nanji_test_predictions.csv", parse_dates=["datetime"])
    monthly = (
        test_df.assign(month=test_df["datetime"].dt.month)
        .groupby("month", as_index=False)[
            ["estimated_active_cars_change", "weighted_extended_prediction", "weather_only_final_prediction"]
        ]
        .mean()
    )
    fig, ax = plt.subplots(figsize=(7.2, 3.7), dpi=180)
    ax.plot(monthly["month"], monthly["estimated_active_cars_change"], color="#0f172a", linewidth=2.6, marker="o", label="Actual")
    ax.plot(monthly["month"], monthly["weighted_extended_prediction"], color="#2563eb", linewidth=2.2, marker="o", label="Best model")
    ax.plot(monthly["month"], monthly["weather_only_final_prediction"], color="#14b8a6", linewidth=2.2, marker="o", label="Operational")
    ax.set_title("2025 월별 평균 변화량 비교")
    ax.set_xlabel("Month")
    ax.set_ylabel("Estimated Active Cars Change")
    ax.grid(alpha=0.2)
    ax.set_xticks(range(1, 13))
    ax.legend(frameon=False, loc="upper right")
    plt.tight_layout()
    fig.savefig(MONTHLY_TREND_IMG, bbox_inches="tight")
    plt.close(fig)

    importance = pd.read_csv(OUT_DIR / "nanji_feature_importance.csv")
    top_importance = (
        importance.loc[importance["model_name"] == "weighted_extended"]
        .sort_values("importance_ratio", ascending=True)
        .tail(8)
    )
    fig, ax = plt.subplots(figsize=(6.1, 3.8), dpi=180)
    ax.barh(top_importance["feature"], top_importance["importance_ratio"], color="#1d4ed8")
    ax.set_title("weighted_extended 주요 변수")
    ax.set_xlabel("Importance Ratio")
    ax.grid(axis="x", alpha=0.2)
    plt.tight_layout()
    fig.savefig(IMPORTANCE_IMG, bbox_inches="tight")
    plt.close(fig)


def build_presentation():
    make_charts()

    compare = pd.read_csv(OUT_DIR / "nanji_weather_only_comparison.csv")
    ranked = compare.sort_values("r2", ascending=False).reset_index(drop=True)
    best_row = ranked.iloc[0]
    final_row = compare.loc[compare["model_name"] == "weather_only_extended_final"].iloc[0]
    weather_row = compare.loc[compare["model_name"] == "weather_only_extended"].iloc[0]

    outlier = pd.read_csv(OUT_DIR / "nanji_outlier_rule_summary.csv").iloc[0]
    removed_rows = pd.read_csv(OUT_DIR / "nanji_feature_removed_rows.csv")
    feature_list = pd.read_csv(OUT_DIR / "nanji_final_operational_feature_list.csv")
    special_cases = pd.read_csv(OUT_DIR / "nanji_special_cases.csv")
    importance = pd.read_csv(OUT_DIR / "nanji_feature_importance.csv")
    weights = pd.read_csv(OUT_DIR / "nanji_weight_table.csv")
    month_df = weights.loc[weights["weight_type"] == "month_weight"].copy()
    month_df["key"] = month_df["key"].astype(int)
    hour_df = weights.loc[weights["weight_type"] == "hour_weight_raw"].copy()
    hour_df["key"] = hour_df["key"].astype(int)
    top_months = ", ".join(str(v) for v in month_df.sort_values("value", ascending=False).head(3)["key"].tolist())
    top_hours = ", ".join(f"{v}시" for v in hour_df.sort_values("value", ascending=False).head(3)["key"].tolist())

    weighted_top = (
        importance.loc[importance["model_name"] == "weighted_extended"]
        .sort_values("importance_ratio", ascending=False)
        .head(5)["feature"]
        .tolist()
    )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_box(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE, WHITE, rounded=False)
    add_box(slide, Inches(0.68), Inches(0.78), Inches(5.6), Inches(5.9), SKY, SKY)
    add_box(slide, Inches(6.62), 0, Inches(6.72), prs.slide_height, NAVY, NAVY, rounded=False)
    add_text(slide, Inches(0.98), Inches(1.18), Inches(4.9), Inches(1.15),
             "난지 ML Change\n가로형 발표자료", size=29, bold=True, color=NAVY)
    add_text(slide, Inches(1.0), Inches(2.42), Inches(4.9), Inches(0.95),
             "기존 점유 차량 수 예측이 아니라\n시간대별 변화량 예측 결과를 기준으로 재구성", size=18, color=TEXT)
    add_bullets(slide, Inches(1.0), Inches(3.58), Inches(4.9), Inches(2.15), [
        "타깃: estimated_active_cars_change",
        f"최고 성능: {best_row['model_name']} / R² {best_row['r2']:.3f}",
        f"운영형 단순안: weather_only_extended_final / 10개 feature",
    ], size=16)
    add_text(slide, Inches(7.15), Inches(0.95), Inches(5.0), Inches(0.45), "이번 보완 포인트", size=24, bold=True, color=WHITE)
    add_box(slide, Inches(7.05), Inches(1.78), Inches(5.15), Inches(1.08), RGBColor(30, 41, 59), RGBColor(30, 41, 59))
    add_text(slide, Inches(7.35), Inches(2.05), Inches(4.55), Inches(0.52),
             "변화량 예측과 점유량 예측을 명확히 구분", size=17, color=WHITE)
    add_box(slide, Inches(7.05), Inches(3.08), Inches(5.15), Inches(1.08), RGBColor(30, 41, 59), RGBColor(30, 41, 59))
    add_text(slide, Inches(7.35), Inches(3.35), Inches(4.55), Inches(0.52),
             "최고 성능 모델과 운영형 모델의 trade-off 반영", size=17, color=WHITE)
    add_box(slide, Inches(7.05), Inches(4.38), Inches(5.15), Inches(1.08), RGBColor(30, 41, 59), RGBColor(30, 41, 59))
    add_text(slide, Inches(7.35), Inches(4.65), Inches(4.6), Inches(0.52),
             "가중치 패턴과 변화량 활용식까지 한 슬라이드 흐름으로 정리", size=17, color=WHITE)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "1. 분석 목적과 데이터 범위", "변화량 예측 기준으로 슬라이드 메시지를 다시 정의")
    add_box(slide, Inches(0.6), Inches(1.35), Inches(3.7), Inches(5.7), LIGHT, BORDER)
    add_text(slide, Inches(0.85), Inches(1.62), Inches(3.1), Inches(0.35), "분석 목표", size=19, bold=True, color=NAVY)
    add_bullets(slide, Inches(0.85), Inches(2.05), Inches(3.05), Inches(4.4), [
        "예측 대상은 특정 시각 점유 차량 수가 아니라 직전 대비 변화량",
        "양수는 유입 우세, 음수는 이탈 우세 흐름으로 해석",
        "실제 운영에서는 현재 차량 수에 변화량 예측치를 더해 다음 시간 점유를 추정",
    ], size=15)
    add_box(slide, Inches(4.5), Inches(1.35), Inches(3.7), Inches(5.7), LIGHT, BORDER)
    add_text(slide, Inches(4.75), Inches(1.62), Inches(3.0), Inches(0.35), "데이터 범위", size=19, bold=True, color=NAVY)
    add_bullets(slide, Inches(4.75), Inches(2.05), Inches(3.05), Inches(4.4), [
        "기간: 2022-01-01 ~ 2025-12-31",
        "총 34,920행, 운영시간 기준 06~23시 중심 예측",
        "train: 2022~2023 / valid: 2024 / test: 2025",
    ], size=15)
    add_box(slide, Inches(8.4), Inches(1.35), Inches(4.3), Inches(5.7), LIGHT, BORDER)
    add_text(slide, Inches(8.65), Inches(1.62), Inches(3.5), Inches(0.35), "활용 방식", size=19, bold=True, color=NAVY)
    add_text(slide, Inches(8.8), Inches(2.18), Inches(3.0), Inches(0.55),
             "next_estimated_active_cars =\ncurrent_actual_cars + predicted_change",
             size=17, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_bullets(slide, Inches(8.7), Inches(3.2), Inches(3.3), Inches(2.6), [
        "변화량 예측이라 TOTAL_CAPACITY만으로 여유면수를 바로 계산하지 않음",
        "현재 점유 상태와 결합해야 다음 시간 점유량이 완성됨",
    ], size=15)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "2. 모델 구조와 품질 점검", "패턴을 먼저 만들고 Ridge가 변화량 오차를 보정")
    steps = [
        ("기본 패턴", "day_type x hour 평균으로\n변화량 base_value 생성"),
        ("월 보정", "month_weight로\n계절별 수준 차이 반영"),
        ("시간 보정", "hour_weight로\n운영시간 내 피크 반영"),
        ("Ridge", "달력, 날씨, 가용성 변수로\n남는 변화량 오차 보정"),
    ]
    for idx, (title, body) in enumerate(steps):
        left = 0.6 + idx * 3.02
        add_box(slide, Inches(left), Inches(1.55), Inches(2.5), Inches(2.2), SKY if idx % 2 == 0 else WHITE, BORDER)
        add_text(slide, Inches(left + 0.2), Inches(1.84), Inches(2.0), Inches(0.35), title, size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(left + 0.18), Inches(2.34), Inches(2.05), Inches(0.85), body, size=14, color=TEXT, align=PP_ALIGN.CENTER)
    add_box(slide, Inches(0.75), Inches(4.25), Inches(3.7), Inches(2.2), LIGHT, BORDER)
    add_text(slide, Inches(1.0), Inches(4.52), Inches(3.1), Inches(0.3), "타깃 이상값 판단", size=18, bold=True, color=NAVY)
    add_bullets(slide, Inches(1.0), Inches(4.95), Inches(3.0), Inches(1.15), [
        f"후보 {int(outlier['candidate_count']):,}건은 모두 특이치로 유지",
        f"실제 제거된 target outlier는 {int(outlier['outlier_count'])}건",
    ], size=14)
    add_box(slide, Inches(4.75), Inches(4.25), Inches(3.9), Inches(2.2), LIGHT, BORDER)
    add_text(slide, Inches(5.0), Inches(4.52), Inches(3.2), Inches(0.3), "입력 변수 정리", size=18, bold=True, color=NAVY)
    add_bullets(slide, Inches(5.0), Inches(4.95), Inches(3.15), Inches(1.15), [
        f"자전거 관련 극단치 {len(removed_rows):,}건만 제거",
        f"특이 고변화 사례 {len(special_cases):,}건은 학습에 유지",
    ], size=14)
    add_box(slide, Inches(8.95), Inches(4.25), Inches(3.7), Inches(2.2), LIGHT, BORDER)
    add_text(slide, Inches(9.2), Inches(4.52), Inches(3.0), Inches(0.3), "핵심 해석", size=18, bold=True, color=NAVY)
    add_bullets(slide, Inches(9.2), Inches(4.95), Inches(2.95), Inches(1.15), [
        "변화량은 0 근처가 많고 큰 증감은 드문 구조",
        "그래서 극단 구간보다 일반 구간 안정성이 더 중요",
    ], size=14)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "3. 패턴 가중치와 변화량 해석", "어느 월, 어느 시간이 변화량 상승 압력을 만드는지")
    slide.shapes.add_picture(str(WEIGHT_IMG), Inches(0.65), Inches(1.35), width=Inches(6.0))
    add_box(slide, Inches(7.0), Inches(1.45), Inches(5.05), Inches(2.0), LIGHT, BORDER)
    add_text(slide, Inches(7.25), Inches(1.72), Inches(4.45), Inches(0.3), "패턴 요약", size=18, bold=True, color=NAVY)
    add_bullets(slide, Inches(7.25), Inches(2.15), Inches(4.35), Inches(1.0), [
        f"월 가중치 상위: {top_months}월",
        f"시간 가중치 상위: {top_hours}",
        "봄·초여름과 초가을, 그리고 점심 직후 시간대에 변화량 상승 압력이 큼",
    ], size=14)
    add_box(slide, Inches(7.0), Inches(3.78), Inches(5.05), Inches(2.2), SOFT, BLUE)
    add_text(slide, Inches(7.25), Inches(4.05), Inches(4.45), Inches(0.3), "발표 멘트용 해석", size=18, bold=True, color=NAVY)
    add_text(slide, Inches(7.28), Inches(4.48), Inches(4.35), Inches(1.05),
             "난지 수요는 단순히 날씨 하나로 움직이지 않고,\n월별 시즌성과 시간대 패턴이 먼저 큰 뼈대를 만들고\n외부 조건이 그 위에 변화를 더하는 구조입니다.",
             size=15, color=TEXT, align=PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "4. 모델 비교와 최종 선택", "최고 성능 모델과 운영형 단순안은 분리해서 설명해야 함")
    slide.shapes.add_picture(str(MODEL_COMPARE_IMG), Inches(0.65), Inches(1.45), width=Inches(5.35))
    table = slide.shapes.add_table(6, 5, Inches(6.15), Inches(1.52), Inches(5.0), Inches(2.85)).table
    headers = ["모델", "alpha", "RMSE", "MAE", "R²"]
    for c, header in enumerate(headers):
        table.cell(0, c).text = header
        style_cell(table.cell(0, c), SKY, NAVY, size=10, bold=True)
    for r in range(len(ranked)):
        row = ranked.iloc[r]
        values = [row["model_name"], f"{row['alpha']:.3f}" if row["alpha"] < 1 else f"{int(row['alpha'])}",
                  f"{row['rmse']:.2f}", f"{row['mae']:.2f}", f"{row['r2']:.3f}"]
        for c, value in enumerate(values):
            table.cell(r + 1, c).text = value
            highlight = row["model_name"] in {"weighted_extended", "weather_only_extended_final"}
            fill = SOFT if highlight else WHITE
            bold = row["model_name"] == "weighted_extended"
            style_cell(table.cell(r + 1, c), fill, TEXT, size=9.2, bold=bold)
    add_box(slide, Inches(6.15), Inches(4.65), Inches(5.0), Inches(1.7), SOFT, BLUE)
    add_bullets(slide, Inches(6.38), Inches(4.92), Inches(4.45), Inches(1.15), [
        f"최고 성능: weighted_extended / R² {best_row['r2']:.3f}, RMSE {best_row['rmse']:.2f}",
        f"운영형 단순안: weather_only_extended_final / R² {final_row['r2']:.3f}",
        f"성능 차이: R² {best_row['r2'] - final_row['r2']:.3f}p, RMSE {final_row['rmse'] - best_row['rmse']:.2f}",
    ], size=14)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "5. 운영형 feature와 중요 변수", "실제 운영에는 10개 feature 단순안이 더 설명하기 쉬움")
    slide.shapes.add_picture(str(FEATURE_GROUP_IMG), Inches(0.7), Inches(1.5), width=Inches(4.7))
    slide.shapes.add_picture(str(IMPORTANCE_IMG), Inches(5.35), Inches(1.5), width=Inches(5.45))
    add_box(slide, Inches(0.85), Inches(5.7), Inches(10.0), Inches(0.9), LIGHT, BORDER)
    add_text(slide, Inches(1.05), Inches(5.95), Inches(9.6), Inches(0.34),
             "운영형 최종안 feature: pattern_prior, month_weight, hour_weight, day_type_offday, hour_sin, hour_cos, month_sin, month_cos, is_holiday, wind_gusts_10m",
             size=14, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(5.6), Inches(6.45), Inches(5.0), Inches(0.28),
             f"최고 성능 모델 핵심 변수: {', '.join(weighted_top[:5])}", size=11, color=MUTED, align=PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "6. 테스트 해석과 적용 방식", "월별 추세는 설명 가능하지만 큰 증감 구간은 보수적으로 봐야 함")
    slide.shapes.add_picture(str(MONTHLY_TREND_IMG), Inches(0.65), Inches(1.38), width=Inches(5.45))
    add_box(slide, Inches(6.25), Inches(1.38), Inches(4.85), Inches(2.1), LIGHT, BORDER)
    add_text(slide, Inches(6.5), Inches(1.65), Inches(4.2), Inches(0.3), "테스트셋 해석", size=18, bold=True, color=NAVY)
    add_bullets(slide, Inches(6.5), Inches(2.05), Inches(4.15), Inches(1.15), [
        f"weighted_extended가 weather_only_extended 대비도 근소 우위(R² {best_row['r2']:.3f} vs {weather_row['r2']:.3f})",
        "9월 변화량 오차가 가장 크고, 20~22시와 16시 구간이 상대적으로 어려움",
    ], size=14)
    add_box(slide, Inches(6.25), Inches(3.78), Inches(4.85), Inches(2.35), SOFT, BLUE)
    add_text(slide, Inches(6.5), Inches(4.05), Inches(4.15), Inches(0.3), "운영 적용 식", size=18, bold=True, color=NAVY)
    add_text(slide, Inches(6.48), Inches(4.48), Inches(4.15), Inches(0.45),
             "next_estimated_active_cars = current_actual_cars + predicted_change",
             size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_bullets(slide, Inches(6.5), Inches(5.05), Inches(4.1), Inches(0.9), [
        "급격한 유입/이탈(절대값 50 이상) 구간은 오차가 커 현장 모니터링 병행 권장",
        "평시 운영 안내와 다음 시간대 흐름 판단에는 충분히 활용 가능",
    ], size=14)

    prs.save(PPT_PATH)
    return PPT_PATH


if __name__ == "__main__":
    output = build_presentation()
    print(output)
