from pathlib import Path

import matplotlib.pyplot as plt
import pandas as pd
from matplotlib import font_manager, rcParams
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
ROOT_DIR = BASE_DIR.parents[1]
SOURCE_DIR = ROOT_DIR / "hmw" / "Note"
OUT_DIR = SOURCE_DIR / "nanji_outputs"

SUMMARY_PATH = BASE_DIR / "hmw_최신보고서_요약.md"
PPT_PATH = BASE_DIR / "hmw_최신보고서_발표자료.pptx"
MODEL_COMPARE_IMG = BASE_DIR / "hmw_model_compare.png"
FEATURE_GROUP_IMG = BASE_DIR / "hmw_feature_groups.png"

MONTH_WEIGHT_IMG = OUT_DIR / "nanji_month_weights.png"
HOUR_WEIGHT_IMG = OUT_DIR / "nanji_hour_weights.png"
MONTHLY_COMPARE_IMG = OUT_DIR / "nanji_2025_monthly_compare.png"

NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
SKY = RGBColor(219, 234, 254)
TEAL = RGBColor(14, 116, 144)
GREEN = RGBColor(22, 163, 74)
ORANGE = RGBColor(234, 88, 12)
RED = RGBColor(220, 38, 38)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(100, 116, 139)
BORDER = RGBColor(203, 213, 225)
LIGHT = RGBColor(248, 250, 252)


def set_font():
    candidates = ["Malgun Gothic", "Apple SD Gothic Neo", "AppleGothic", "NanumGothic"]
    available = {f.name for f in font_manager.fontManager.ttflist}
    for candidate in candidates:
        if candidate in available:
            rcParams["font.family"] = candidate
            return candidate
    return "Malgun Gothic"


FONT_NAME = set_font()


def add_box(slide, left, top, width, height, fill, line=None, rounded=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_text(slide, left, top, width, height, text, size=20, bold=False,
             color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(5)
        p.line_spacing = 1.15
        run = p.runs[0]
        run.font.name = FONT_NAME
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_header(slide, title, subtitle=""):
    add_text(slide, Inches(0.55), Inches(0.3), Inches(12.0), Inches(0.45), title, size=24, bold=True, color=NAVY)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.82), Inches(12.0), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.color.rgb = BLUE
    if subtitle:
        add_text(slide, Inches(0.57), Inches(0.9), Inches(11.8), Inches(0.25), subtitle, size=10, color=MUTED)


def style_cell(cell, fill, font_color, size=10, bold=False):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = FONT_NAME
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = font_color


def make_charts():
    compare = pd.read_csv(OUT_DIR / "nanji_weather_only_comparison.csv").sort_values("r2", ascending=True)
    colors = ["#cbd5e1", "#93c5fd", "#60a5fa", "#2563eb", "#1d4ed8"]
    fig, ax = plt.subplots(figsize=(7.2, 3.8), dpi=180)
    bars = ax.barh(compare["model_name"], compare["r2"], color=colors[: len(compare)])
    ax.set_xlim(0.72, 0.77)
    ax.set_xlabel("Test R²")
    ax.set_title("운영 후보 모델 성능 비교")
    ax.grid(axis="x", alpha=0.2)
    for bar, val in zip(bars, compare["r2"]):
        ax.text(val + 0.0005, bar.get_y() + bar.get_height() / 2, f"{val:.3f}", va="center", fontsize=8)
    plt.tight_layout()
    fig.savefig(MODEL_COMPARE_IMG, bbox_inches="tight")
    plt.close(fig)

    features = pd.read_csv(OUT_DIR / "nanji_final_operational_feature_list.csv")
    group_counts = features.groupby("feature_group").size().sort_values(ascending=True)
    fig, ax = plt.subplots(figsize=(5.8, 3.4), dpi=180)
    bars = ax.barh(group_counts.index, group_counts.values, color=["#dbeafe", "#93c5fd", "#38bdf8", "#0ea5e9"])
    ax.set_xlabel("Count")
    ax.set_title("최종 운영형 feature 구성")
    ax.grid(axis="x", alpha=0.2)
    for bar, val in zip(bars, group_counts.values):
        ax.text(val + 0.03, bar.get_y() + bar.get_height() / 2, str(val), va="center", fontsize=8)
    plt.tight_layout()
    fig.savefig(FEATURE_GROUP_IMG, bbox_inches="tight")
    plt.close(fig)


def build_presentation():
    make_charts()

    compare = pd.read_csv(OUT_DIR / "nanji_weather_only_comparison.csv")
    final_row = compare.loc[compare["model_name"] == "weather_only_extended_final"].iloc[0]
    weighted_row = compare.loc[compare["model_name"] == "weighted_extended"].iloc[0]
    core_row = compare.loc[compare["model_name"] == "weighted_core"].iloc[0]

    outlier = pd.read_csv(OUT_DIR / "nanji_outlier_rule_summary.csv").iloc[0]
    feature_pruning = pd.read_csv(OUT_DIR / "nanji_feature_pruning.csv")
    removed_rows = pd.read_csv(OUT_DIR / "nanji_feature_removed_rows.csv")
    feature_list = pd.read_csv(OUT_DIR / "nanji_final_operational_feature_list.csv")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_box(slide, 0, 0, prs.slide_width, prs.slide_height, WHITE, WHITE, rounded=False)
    add_box(slide, Inches(0.65), Inches(0.8), Inches(5.5), Inches(5.7), SKY, SKY)
    add_box(slide, Inches(6.6), Inches(0), Inches(6.75), prs.slide_height, NAVY, NAVY, rounded=False)
    add_text(slide, Inches(0.95), Inches(1.2), Inches(4.8), Inches(1.2),
             "HMW 최신 보고서\n발표자료", size=28, bold=True, color=NAVY)
    add_text(slide, Inches(0.98), Inches(2.35), Inches(4.9), Inches(1.0),
             "난지 한강공원 시간대별 주차 수요 예측 모델을\n발표용 흐름으로 다시 정리한 버전", size=18, color=TEXT)
    add_bullets(slide, Inches(1.0), Inches(3.55), Inches(4.8), Inches(2.0), [
        "기준 문서: nanji_weighted_ridge_modeling_report.md",
        "핵심 주제: 운영 가능한 주차 수요 예측 모델 선정",
        f"최종 운영형 성능: R² {final_row['r2']:.3f}, RMSE {final_row['rmse']:.2f}",
    ], size=16)
    add_text(slide, Inches(7.15), Inches(1.0), Inches(4.9), Inches(0.5), "발표 핵심", size=24, bold=True, color=WHITE)
    add_box(slide, Inches(7.1), Inches(1.85), Inches(5.1), Inches(1.0), RGBColor(30, 41, 59), RGBColor(30, 41, 59))
    add_text(slide, Inches(7.35), Inches(2.05), Inches(4.5), Inches(0.5),
             "직접 로그가 아닌 추정형 시간별 데이터셋 기반", size=17, color=WHITE)
    add_box(slide, Inches(7.1), Inches(3.1), Inches(5.1), Inches(1.0), RGBColor(30, 41, 59), RGBColor(30, 41, 59))
    add_text(slide, Inches(7.35), Inches(3.3), Inches(4.5), Inches(0.5),
             "패턴 + 가중치 + Ridge 구조로 해석력 확보", size=17, color=WHITE)
    add_box(slide, Inches(7.1), Inches(4.35), Inches(5.1), Inches(1.0), RGBColor(30, 41, 59), RGBColor(30, 41, 59))
    add_text(slide, Inches(7.35), Inches(4.55), Inches(4.6), Inches(0.5),
             "운영용 최종안은 13개 feature의 날씨 기반 모델", size=17, color=WHITE)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "1. 분석 목적과 데이터 범위", "무엇을 예측했고 어떤 데이터로 학습했는지")
    add_box(slide, Inches(0.6), Inches(1.35), Inches(3.6), Inches(5.7), LIGHT, BORDER)
    add_text(slide, Inches(0.85), Inches(1.62), Inches(3.0), Inches(0.35), "분석 목표", size=19, bold=True, color=NAVY)
    add_bullets(slide, Inches(0.85), Inches(2.05), Inches(2.95), Inches(4.4), [
        "예측 대상은 특정 시각의 추정 주차 차량 수인 estimated_active_cars",
        "최종적으로는 여유 주차공간 수로 바꿔 운영에 활용하는 것이 목적",
        "해석 가능성과 실제 운영 연결성을 함께 고려",
    ], size=15)
    add_box(slide, Inches(4.45), Inches(1.35), Inches(3.8), Inches(5.7), LIGHT, BORDER)
    add_text(slide, Inches(4.7), Inches(1.62), Inches(3.0), Inches(0.35), "데이터 범위", size=19, bold=True, color=NAVY)
    add_bullets(slide, Inches(4.7), Inches(2.05), Inches(3.1), Inches(4.4), [
        "기간: 2022-01-01 ~ 2025-12-31",
        "총 34,920행, 중복 0, 타깃 결측 0",
        "직접 실측 시간 로그가 아니라 일별 원본을 24시간으로 확장한 추정형 데이터셋",
    ], size=15)
    add_box(slide, Inches(8.5), Inches(1.35), Inches(4.2), Inches(5.7), LIGHT, BORDER)
    add_text(slide, Inches(8.75), Inches(1.62), Inches(3.5), Inches(0.35), "운영 조건과 분할", size=19, bold=True, color=NAVY)
    add_bullets(slide, Inches(8.75), Inches(2.05), Inches(3.3), Inches(4.4), [
        "운영시간은 06~23시로 간주",
        "00~05시는 운영 외 구간으로 보고 0 유지",
        "train: 2022~2023 / valid: 2024 / test: 2025",
    ], size=15)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "2. 전처리와 모델 구조", "패턴을 먼저 설명하고 남는 차이를 Ridge가 보정")
    add_box(slide, Inches(0.65), Inches(1.5), Inches(2.45), Inches(2.05), SKY, BLUE)
    add_text(slide, Inches(0.92), Inches(1.82), Inches(1.9), Inches(0.35), "1. 기본 패턴", size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.9), Inches(2.25), Inches(1.95), Inches(0.7),
             "day_type x hour 평균으로\nbase_value 생성", size=15, color=TEXT, align=PP_ALIGN.CENTER)
    add_box(slide, Inches(3.35), Inches(1.5), Inches(2.45), Inches(2.05), LIGHT, BORDER)
    add_text(slide, Inches(3.62), Inches(1.82), Inches(1.9), Inches(0.35), "2. 월 보정", size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(3.6), Inches(2.25), Inches(1.95), Inches(0.7),
             "month_weight로\n계절별 수준 차이 반영", size=15, color=TEXT, align=PP_ALIGN.CENTER)
    add_box(slide, Inches(6.05), Inches(1.5), Inches(2.45), Inches(2.05), SKY, BLUE)
    add_text(slide, Inches(6.32), Inches(1.82), Inches(1.9), Inches(0.35), "3. 시간 보정", size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(6.3), Inches(2.25), Inches(1.95), Inches(0.7),
             "hour_weight로\n운영시간대 혼잡 차이 반영", size=15, color=TEXT, align=PP_ALIGN.CENTER)
    add_box(slide, Inches(8.75), Inches(1.5), Inches(2.45), Inches(2.05), LIGHT, BORDER)
    add_text(slide, Inches(9.02), Inches(1.82), Inches(1.9), Inches(0.35), "4. Ridge 보정", size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(9.0), Inches(2.25), Inches(1.95), Inches(0.7),
             "날씨·주기 변수와 함께\n최종 오차 축소", size=15, color=TEXT, align=PP_ALIGN.CENTER)
    add_box(slide, Inches(0.7), Inches(4.15), Inches(4.2), Inches(2.4), LIGHT, BORDER)
    add_text(slide, Inches(0.95), Inches(4.42), Inches(3.5), Inches(0.35), "전처리 핵심 결과", size=18, bold=True, color=NAVY)
    add_bullets(slide, Inches(0.95), Inches(4.82), Inches(3.5), Inches(1.4), [
        f"타깃 상한 초과 후보 {int(outlier['candidate_count']):,}건은 모두 특이값으로 유지",
        f"입력 변수 쪽 실제 제거 행은 총 {len(removed_rows):,}건",
        f"상관계수 0.9 이상 기준으로 {len(feature_pruning):,}개 중복 관계 정리",
    ], size=14)
    add_box(slide, Inches(5.2), Inches(4.15), Inches(7.45), Inches(2.4), LIGHT, BORDER)
    add_text(slide, Inches(5.45), Inches(4.42), Inches(6.8), Inches(0.35), "해석 포인트", size=18, bold=True, color=NAVY)
    add_bullets(slide, Inches(5.45), Inches(4.82), Inches(6.7), Inches(1.4), [
        "작은 값과 중간 값이 많고 큰 값은 드문 long-tail 분포",
        "그래서 전체를 한 번에 맞추는 것보다 패턴을 먼저 설명하는 구조가 더 안정적",
        "sin/cos 변수로 시간과 계절의 주기성을 추가 반영",
    ], size=14)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "3. 성능 비교와 최종 운영형 모델 선정", "좋은 성능과 미래 입력 확보 가능성을 함께 고려")
    slide.shapes.add_picture(str(MODEL_COMPARE_IMG), Inches(0.7), Inches(1.45), width=Inches(5.8))
    add_box(slide, Inches(6.85), Inches(1.45), Inches(5.85), Inches(2.45), LIGHT, BORDER)
    add_text(slide, Inches(7.12), Inches(1.72), Inches(5.2), Inches(0.35), "핵심 비교 결과", size=18, bold=True, color=NAVY)
    add_bullets(slide, Inches(7.12), Inches(2.12), Inches(5.1), Inches(1.4), [
        f"weighted_core: R² {core_row['r2']:.3f}, RMSE {core_row['rmse']:.2f}",
        f"weighted_extended: R² {weighted_row['r2']:.3f}, RMSE {weighted_row['rmse']:.2f}",
        f"weather_only_extended_final: R² {final_row['r2']:.3f}, RMSE {final_row['rmse']:.2f}",
    ], size=14)
    add_box(slide, Inches(6.85), Inches(4.1), Inches(5.85), Inches(2.2), SKY, BLUE)
    add_text(slide, Inches(7.12), Inches(4.38), Inches(5.2), Inches(0.35), "최종 선정 이유", size=18, bold=True, color=NAVY)
    add_bullets(slide, Inches(7.12), Inches(4.78), Inches(5.15), Inches(1.25), [
        "교통·자전거·행사 변수는 미래 시점 실사용이 어려움",
        "날씨는 예보 데이터를 붙일 수 있어 운영 연결성이 높음",
        "성능도 가장 좋고 입력 구조도 단순해 최종 운영안으로 적합",
    ], size=14)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "4. 최종 운영형 feature 구성", "13개 feature로 단순화해 운영 가능성 확보")
    slide.shapes.add_picture(str(FEATURE_GROUP_IMG), Inches(0.8), Inches(1.55), width=Inches(4.9))
    add_box(slide, Inches(6.1), Inches(1.5), Inches(6.0), Inches(4.6), LIGHT, BORDER)
    add_text(slide, Inches(6.35), Inches(1.78), Inches(5.2), Inches(0.35), "최종 feature", size=18, bold=True, color=NAVY)
    grouped = {
        "패턴": feature_list.loc[feature_list["feature_group"] == "pattern", "feature"].tolist(),
        "주기": feature_list.loc[feature_list["feature_group"] == "cycle", "feature"].tolist(),
        "달력": feature_list.loc[feature_list["feature_group"] == "calendar", "feature"].tolist(),
        "날씨": feature_list.loc[feature_list["feature_group"] == "weather", "feature"].tolist(),
    }
    y = 2.18
    for label, values in grouped.items():
        add_text(slide, Inches(6.35), Inches(y), Inches(0.8), Inches(0.25), f"{label}", size=14, bold=True, color=BLUE)
        add_text(slide, Inches(7.2), Inches(y), Inches(4.5), Inches(0.4), ", ".join(values), size=13, color=TEXT)
        y += 0.68
    add_box(slide, Inches(0.8), Inches(6.2), Inches(11.3), Inches(0.8), SKY, BLUE)
    add_text(slide, Inches(1.0), Inches(6.42), Inches(10.8), Inches(0.3),
             "정리하면 패턴 뼈대는 유지하고, 미래 시점에서 바로 확보 가능한 날씨 변수만 남긴 운영형 Ridge 모델입니다.",
             size=15, color=NAVY, align=PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "5. 계절·시간 패턴 시각화", "모델의 큰 뼈대는 month_weight와 hour_weight가 설명")
    if MONTH_WEIGHT_IMG.exists():
        slide.shapes.add_picture(str(MONTH_WEIGHT_IMG), Inches(0.7), Inches(1.45), width=Inches(5.85))
    if HOUR_WEIGHT_IMG.exists():
        slide.shapes.add_picture(str(HOUR_WEIGHT_IMG), Inches(6.8), Inches(1.45), width=Inches(5.85))
    add_box(slide, Inches(0.9), Inches(6.0), Inches(11.0), Inches(0.7), LIGHT, BORDER)
    add_text(slide, Inches(1.1), Inches(6.18), Inches(10.6), Inches(0.25),
             "월별 수요 수준 차이와 운영시간대 혼잡 차이가 최종 예측의 핵심 축이며, 이후 날씨가 추가 설명력을 더하는 구조입니다.",
             size=15, color=TEXT, align=PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "6. 모델 한계와 운영 해석", "일반 구간에서는 실용적이지만 고혼잡 구간은 보수적으로 봐야 함")
    add_box(slide, Inches(0.7), Inches(1.45), Inches(5.0), Inches(4.9), LIGHT, BORDER)
    add_text(slide, Inches(0.98), Inches(1.72), Inches(4.3), Inches(0.35), "모델 한계", size=18, bold=True, color=NAVY)
    add_bullets(slide, Inches(0.98), Inches(2.12), Inches(4.2), Inches(3.6), [
        "실제값 500 이상 고혼잡 구간에서는 과소예측 경향이 나타남",
        "오차가 큰 시간대는 주로 12~19시, 특히 13~16시",
        "주말·휴일 성격이 강한 날과 4~6월, 9~10월이 상대적으로 어려운 구간",
    ], size=15)
    add_box(slide, Inches(6.0), Inches(1.45), Inches(6.0), Inches(4.9), LIGHT, BORDER)
    add_text(slide, Inches(6.28), Inches(1.72), Inches(5.2), Inches(0.35), "운영 적용 방식", size=18, bold=True, color=NAVY)
    add_bullets(slide, Inches(6.28), Inches(2.12), Inches(5.2), Inches(3.6), [
        "1시간 후, 2시간 후, 하루 뒤 같은 단기 예측에는 활용 가능",
        "예측된 estimated_active_cars를 총 주차면 수와 비교해 여유 주차공간 수로 변환",
        "피크 시간대는 과소예측 가능성을 함께 안내하는 방식이 적절",
    ], size=15)
    if MONTHLY_COMPARE_IMG.exists():
        slide.shapes.add_picture(str(MONTHLY_COMPARE_IMG), Inches(8.15), Inches(4.35), width=Inches(3.4))

    slide = prs.slides.add_slide(blank)
    add_header(slide, "7. 결론", "발표 마무리용 한 장 요약")
    add_box(slide, Inches(0.85), Inches(1.45), Inches(11.5), Inches(4.9), SKY, BLUE)
    add_bullets(slide, Inches(1.2), Inches(1.95), Inches(10.8), Inches(3.9), [
        "난지 주차 수요 예측은 직접 로그가 아닌 추정형 시간별 데이터셋 기반으로 수행했다.",
        "모델 구조는 기본 패턴 + month/hour 가중치 + Ridge 회귀로 해석력을 확보했다.",
        "최종 운영안은 13개 feature를 쓰는 weather_only_extended_final 모델이다.",
        f"최종 test 성능은 R² {final_row['r2']:.3f}, RMSE {final_row['rmse']:.2f}, MAE {final_row['mae']:.2f}다.",
        "일반 운영 구간에서는 실용적이지만, 500 이상 고혼잡 구간은 보수적으로 해석해야 한다.",
        "실무에서는 예측 점유량을 총 주차면 수와 비교해 여유 주차공간 수로 바꿔 사용하는 방식이 적절하다.",
    ], size=18)
    add_text(slide, Inches(0.95), Inches(6.7), Inches(10.5), Inches(0.25),
             "생성 파일: yeeun/0412/hmw_최신보고서_발표자료.pptx", size=10, color=MUTED)

    prs.save(PPT_PATH)


if __name__ == "__main__":
    build_presentation()
